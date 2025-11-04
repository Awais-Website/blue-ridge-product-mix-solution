"""
create_ppt.py
==============

This script uses the `python-pptx` library to assemble a short slide deck
for the Blue Ridge Hot Tubs product‑mix and sensitivity analysis.  The
resulting presentation illustrates the business problem, the LP model,
baseline solution and the findings from the sensitivity analysis.

Run this script to generate a PowerPoint file named
`Blue_Ridge_Sensitivity_Analysis.pptx` in the same directory.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def add_title_slide(prs):
    slide_layout = prs.slide_layouts[0]  # title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Blue Ridge Hot Tubs – Product Mix & Sensitivity Analysis"
    subtitle.text = "Using Python LP Optimisation"


def add_problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # title & content
    slide.shapes.title.text = "Business Problem"
    body = slide.shapes.placeholders[1].text_frame
    body.text = "Blue Ridge Hot Tubs manufactures two models of hot tubs."
    p = body.add_paragraph()
    p.text = "• Aqua-Spas: $350 profit; uses 1 pump, 9 labour hours, 12 feet tubing"
    p.level = 0
    p = body.add_paragraph()
    p.text = "• Hydro-Luxes: $300 profit; uses 1 pump, 6 labour hours, 16 feet tubing"
    p.level = 0
    p = body.add_paragraph()
    p.text = "Available resources: 200 pumps, 1,566 labour hours, 2,880 feet of tubing"
    p.level = 0
    p = body.add_paragraph()
    p.text = "Goal: maximise total profit while respecting resource constraints"
    p.level = 0


def add_model_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "LP Model"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Decision variables:"
    p = tf.add_paragraph()
    p.text = "x₁ = number of Aqua-Spas"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "x₂ = number of Hydro-Luxes"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Objective: maximise 350 x₁ + 300 x₂"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Subject to:"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "• x₁ + x₂ ≤ 200  (Pumps)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 9x₁ + 6x₂ ≤ 1,566  (Labour hours)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 12x₁ + 16x₂ ≤ 2,880  (Tubing feet)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• x₁, x₂ ≥ 0"
    p.level = 1


def add_baseline_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Baseline Solution"
    body = slide.shapes.placeholders[1].text_frame
    body.text = "Optimal production plan (continuous LP):"
    p = body.add_paragraph()
    p.text = "Aqua-Spas = 122 units"
    p.level = 1
    p = body.add_paragraph()
    p.text = "Hydro-Luxes = 78 units"
    p.level = 1
    p = body.add_paragraph()
    p.text = "Maximum profit = $66,100"
    p.level = 1
    p = body.add_paragraph()
    p.text = "Resources used: all pumps and labour; tubing has slack"
    p.level = 1


def add_sensitivity_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Sensitivity Analysis Results"
    body = slide.shapes.placeholders[1].text_frame
    body.text = "Marginal values (constant region):"
    p = body.add_paragraph()
    p.text = "• Pumps: Profit increases by $200 per pump up to ≈207 pumps"
    p.level = 1
    p = body.add_paragraph()
    p.text = "• Labour: Profit increases by $16.67 per hour up to ≈1800 hours"
    p.level = 1
    p = body.add_paragraph()
    p.text = "• Tubing: Additional tubing does not change profit (non-binding)"
    p.level = 1
    # Insert the sensitivity slide image if present
    img_path = Path(__file__).resolve().parent.parent / "images" / "sensitivity_analysis.png"
    if img_path.exists():
        left = Inches(5.2)
        top = Inches(1.5)
        width = Inches(4.0)
        slide.shapes.add_picture(str(img_path), left, top, width=width)


def add_conclusion_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Conclusions & Recommendations"
    body = slide.shapes.placeholders[1].text_frame
    body.text = "Findings:"
    p = body.add_paragraph()
    p.text = "• Pumps and labour are binding in the baseline solution; tubing has slack"
    p.level = 1
    p = body.add_paragraph()
    p.text = "• Pumps add $200 in profit up to ~207 pumps"
    p.level = 1
    p = body.add_paragraph()
    p.text = "• Labour adds $16.67 per hour up to ~1800 hours"
    p.level = 1
    p = body.add_paragraph()
    p.text = "• Investing in additional tubing is not beneficial in this region"
    p.level = 1
    p = body.add_paragraph()
    p.text = "Recommendation: Prioritise acquiring more pumps and labour up to breakpoints"
    p.level = 0


def create_presentation(output_path: Path) -> None:
    prs = Presentation()
    # Remove default slide if present (python-pptx sometimes inserts one)
    if prs.slides:
        prs.slides.remove(prs.slides[0])
    add_title_slide(prs)
    add_problem_slide(prs)
    add_model_slide(prs)
    add_baseline_slide(prs)
    add_sensitivity_slide(prs)
    add_conclusion_slide(prs)
    prs.save(str(output_path))


if __name__ == "__main__":
    here = Path(__file__).resolve().parent
    out = here / "Blue_Ridge_Sensitivity_Analysis.pptx"
    create_presentation(out)
    print(f"Presentation saved to {out}")